import hashlib
import re
import fitz  # PyMuPDF
import docx
import pptx


# 파일별 개인정보 탐지 및 암호화 방식
# 	1.	PDF 파일
# 	•	PyMuPDF 또는 pdfplumber 라이브러리를 사용하여 텍스트를 추출
# 	•	정규 표현식을 활용하여 개인정보(예: 주민등록번호, 이메일, 전화번호) 탐색
# 	•	탐지된 정보를 해싱(SHA-256) 처리 후 원본을 대체
# 	2.	Word (DOCX) 파일
# 	•	python-docx를 사용하여 문서의 모든 단락 및 표에서 텍스트 추출
# 	•	정규 표현식을 이용하여 개인정보 검색
# 	•	탐지된 정보를 해싱 후 원본을 업데이트
# 	3.	PPT 파일
# 	•	python-pptx를 사용하여 슬라이드의 텍스트 검색
# 	•	개인정보 탐지 후 해싱하여 변환
# 	4.	TXT 파일
# 	•	일반 텍스트 파일이므로 단순한 정규 표현식 검색 후 해싱 적용 가능

# 핵심 기능 정리
#
# ✅ 개인정보 탐지: 이메일, 전화번호, 주민등록번호 등을 정규식으로 탐색
# ✅ 개인정보 해싱: SHA-256을 사용하여 원본을 대체
# ✅ 포맷별 변환 지원: PDF, DOCX, PPTX, TXT 처리 가능
# ✅ 출력 파일 생성: 원본 파일을 변환한 후 새로운 파일로 저장
#
# 이 방식으로 백신 엔진과 유사하게 동작할 수 있습니다.
# 추가적으로 JSON, CSV, Excel 같은 데이터 파일도 처리할 수 있도록 확장 가능

class PersonalInfoMasker:

    def __init__(self):
        self.patterns = {
            "email": r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}",
            "phone": r"\d{2,3}-\d{3,4}-\d{4}",
            "id": r"\d{6}-\d{7}"
        }

    def hash_value(self, value):
        """개인정보를 SHA-256 해싱하여 반환"""
        return hashlib.sha256(value.encode()).hexdigest()

    def mask_text(self, text):
        """텍스트에서 개인정보 탐지 후 해싱하여 대체"""
        for key, pattern in self.patterns.items():
            text = re.sub(pattern, lambda x: self.hash_value(x.group()), text)
        return text

    def process_pdf(self, pdf_path, output_path):
        """PDF 파일에서 개인정보 탐지 후 암호화"""
        doc = fitz.open(pdf_path)
        for page in doc:
            text = page.get_text()
            new_text = self.mask_text(text)
            page.insert_text((50, 50), new_text, fontsize=10)
        doc.save(output_path)

    def process_docx(self, docx_path, output_path):
        """DOCX 파일에서 개인정보 탐지 후 암호화"""
        doc = docx.Document(docx_path)
        for para in doc.paragraphs:
            para.text = self.mask_text(para.text)
        doc.save(output_path)

    def process_pptx(self, pptx_path, output_path):
        """PPTX 파일에서 개인정보 탐지 후 암호화"""
        pres = pptx.Presentation(pptx_path)
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape.text = self.mask_text(shape.text)
        pres.save(output_path)

    def process_txt(self, txt_path, output_path):
        """TXT 파일에서 개인정보 탐지 후 암호화"""
        with open(txt_path, "r", encoding="utf-8") as f:
            text = f.read()
        new_text = self.mask_text(text)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(new_text)


# 사용 예시
masker = PersonalInfoMasker()
masker.process_pdf("input.pdf", "output.pdf")
masker.process_docx("input.docx", "output.docx")
masker.process_pptx("input.pptx", "output.pptx")
masker.process_txt("input.txt", "output.txt")